from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------- Brand palette ----------
MINT = RGBColor(94, 234, 212)  # #5EEAD4
INK = RGBColor(11, 13, 18)  # #0B0D12
SLATE = RGBColor(15, 23, 42)  # #0F172A
S500 = RGBColor(71, 85, 105)  # slate-500
S600 = RGBColor(100, 116, 139)  # slate-600
WHITE = RGBColor(255, 255, 255)

prs = Presentation()
prs.slide_width = Inches(13.33)  # 1280 x 720 @ 96dpi (16:9)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]


def mint_bar_footer(slide, text="Animica — animica.dev"):
    # Mint footer bar + small label
    bar_h = 0.28
    bar = slide.shapes.add_shape(
        MSO.RECTANGLE,
        Inches(0),
        prs.slide_height - Inches(bar_h),
        prs.slide_width,
        Inches(bar_h),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = MINT
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(
        Inches(0.35),
        prs.slide_height - Inches(bar_h) + Inches(0.06),
        Inches(6),
        Inches(0.4),
    )
    tf = tb.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = INK


def title_shape(slide, text, y=1.0):
    box = slide.shapes.add_textbox(Inches(1.3), Inches(y), Inches(10.7), Inches(1.4))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = SLATE
    return box


def subtitle_shape(slide, text, y=2.2, size=26, color=S500):
    box = slide.shapes.add_textbox(Inches(1.3), Inches(y), Inches(10.7), Inches(0.9))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    return box


def bullets(slide, items, x=1.3, y=3.0, w=10.7, h=3.8, size=26, color=SLATE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {it}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def pill(slide, x=1.3, y=0.6, w=1.7, h=0.5, text="ANIMICA"):
    sh = slide.shapes.add_shape(
        MSO.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = MINT
    sh.line.fill.background()
    tf = sh.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = INK
    p.alignment = PP_ALIGN.CENTER
    return sh


def three_cols(slide, cols, start_y=2.2):
    # cols: list of (title, body)
    col_w = 3.6
    gap = 0.3
    x = 1.3
    for title, body in cols:
        hb = slide.shapes.add_textbox(
            Inches(x), Inches(start_y), Inches(col_w), Inches(0.6)
        )
        ht = hb.text_frame
        ht.clear()
        p = ht.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(28)
        p.font.color.rgb = SLATE

        bb = slide.shapes.add_textbox(
            Inches(x), Inches(start_y + 0.7), Inches(col_w), Inches(2.4)
        )
        bt = bb.text_frame
        bt.clear()
        for i, line in enumerate(body.split("\n")):
            p = bt.paragraphs[0] if i == 0 else bt.add_paragraph()
            p.text = line
            p.font.size = Pt(22)
            p.font.color.rgb = S600
        x += col_w + gap


# ===== Cover =====
s = prs.slides.add_slide(blank)
pill(s)
title_shape(s, "Animica Pitch Deck", y=1.6)
subtitle_shape(
    s, "Deterministic Python VM · Post-Quantum Security · Useful-Work Mining", y=2.6
)
subtitle_shape(s, "v1.0 — {date}".format(date="2025-11-01"), y=3.2, size=16, color=S600)
# Accent ring motif (bottom-left)
ring = s.shapes.add_shape(MSO.OVAL, Inches(0.5), Inches(6.4), Inches(0.6), Inches(0.6))
ring.fill.background()
ring.line.width = Pt(4)
ring.line.color.rgb = MINT
spark = s.shapes.add_shape(
    MSO.OVAL, Inches(0.85), Inches(6.55), Inches(0.14), Inches(0.14)
)
spark.fill.solid()
spark.fill.fore_color.rgb = MINT
spark.line.fill.background()
mint_bar_footer(s)

# ===== Problem =====
s = prs.slides.add_slide(blank)
pill(s)
title_shape(s, "Problem")
bullets(
    s,
    [
        "Smart contracts are non-deterministic across environments, complicating audits.",
        "Post-quantum migration paths are unclear for most L1s.",
        "Upgrade governance is risky without machine-readable registries & guardrails.",
    ],
)
mint_bar_footer(s)

# ===== Solution =====
s = prs.slides.add_slide(blank)
pill(s)
title_shape(s, "Solution")
three_cols(
    s,
    [
        (
            "Deterministic VM",
            "Python-VM with a constrained stdlib\nCanonical hashing & ABI\nGas-metered I/O",
        ),
        (
            "Post-Quantum First",
            "Dilithium3 signatures\nKyber-768 KEX (rotatable)\nRotation policies in governance",
        ),
        (
            "Governance, Formalized",
            "Schemas & registries in-repo\nParam bounds & CI checks\nSafe, staged upgrades",
        ),
    ],
)
mint_bar_footer(s)

# ===== Product =====
s = prs.slides.add_slide(blank)
pill(s)
title_shape(s, "Product")
bullets(
    s,
    [
        "Explorer + Wallet Extension + Flutter Wallet",
        "DEX + CEX reference implementations",
        "Data Availability layer & Randomness beacons",
        "SDKs: TypeScript & Python",
    ],
)
mint_bar_footer(s)

# ===== Traction (placeholders) =====
s = prs.slides.add_slide(blank)
pill(s)
title_shape(s, "Traction")
three_cols(
    s,
    [
        ("Developers", "1.2k+ SDK downloads/mo\n>120 contracts deployed on testnet"),
        ("Network", "Avg. 9s block time\n>99.9% uptime last 90d"),
        ("Ecosystem", "3 launch partners\n2 audits in progress"),
    ],
)
mint_bar_footer(s)

# ===== Market =====
s = prs.slides.add_slide(blank)
pill(s)
title_shape(s, "Market")
three_cols(
    s,
    [
        ("TAM", "$40B+ blockchain infra / yr"),
        ("SAM", "$8B smart-contract platforms"),
        ("SOM", "$500M immediate target"),
    ],
)
subtitle_shape(
    s,
    "Focus: builders & security-sensitive apps (fintech, RWA, compliance).",
    y=4.4,
    size=22,
    color=S600,
)
mint_bar_footer(s)

# ===== Business Model =====
s = prs.slides.add_slide(blank)
pill(s)
title_shape(s, "Business Model")
bullets(
    s,
    [
        "Block rewards + fees (on-chain economy)",
        "Enterprise support & managed nodes",
        "Hosted DA / randomness / AICF credits",
        "Grants & ecosystem fund for strategic growth",
    ],
)
mint_bar_footer(s)

# ===== Technology =====
s = prs.slides.add_slide(blank)
pill(s)
title_shape(s, "Technology")
bullets(
    s,
    [
        "Deterministic Python VM · ABI & gas model",
        "PQ identities: Dilithium3, Kyber-768 (rotatable)",
        "Data Availability: NMT blobs, RS coding",
        "Randomness: VDF + optional QRNG feed",
        "PoIES consensus with fairness caps Γ",
    ],
)
mint_bar_footer(s)

# ===== Roadmap =====
s = prs.slides.add_slide(blank)
pill(s)
title_shape(s, "Roadmap")
bullets(
    s,
    [
        "Q1: Public testnet, SDKs, DEX/CEX refs",
        "Q2: Audits, PQ rotation dry-run, governance boot",
        "Q3: Mainnet candidate, staged rollout, canaries",
        "Q4: Ecosystem programs, AICF integrations",
    ],
)
mint_bar_footer(s)

# ===== Team =====
s = prs.slides.add_slide(blank)
pill(s)
title_shape(s, "Team")
three_cols(
    s,
    [
        ("Founder/CEO", "Protocol, consensus, security"),
        ("Head of Eng", "Runtime, clients, dev-tools"),
        ("Head of Ecosystem", "Partners, grants, growth"),
    ],
)
subtitle_shape(
    s, "Advisors: cryptography, governance, infra.", y=4.4, size=22, color=S600
)
mint_bar_footer(s)

# ===== The Ask =====
s = prs.slides.add_slide(blank)
pill(s)
title_shape(s, "The Ask")
bullets(
    s,
    [
        "Seeking: $X seed to mainnet (18 months runway)",
        "Use of proceeds: protocol hires, audits, validator grants",
        "Contact: founders@animica.dev · animica.dev/press",
    ],
)
mint_bar_footer(s)

out = "contrib/docs-templates/slide-deck/Animica-Pitch.pptx"
prs.save(out)
print("✅ Wrote", out)
