from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Brand colors
MINT = RGBColor(94, 234, 212)
INK = RGBColor(11, 13, 18)
SLATE = RGBColor(15, 23, 42)
SLATE_500 = RGBColor(71, 85, 105)
SLATE_600 = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ===== Cover slide =====
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)
shape = slide.shapes.add_shape(1, Inches(2), Inches(1.5), Inches(2), Inches(0.6))
shape.fill.solid()
shape.fill.fore_color.rgb = MINT
shape.line.fill.background()
tx = shape.text_frame
tx.text = "ANIMICA"
p = tx.paragraphs[0]
p.font.bold = True
p.font.size = Pt(20)
p.font.color.rgb = INK
p.alignment = PP_ALIGN.CENTER

# Title
title_box = slide.shapes.add_textbox(Inches(2), Inches(2.8), Inches(10), Inches(1.5))
tf = title_box.text_frame
tf.text = "Whitepaper Title"
p = tf.paragraphs[0]
p.font.bold = True
p.font.size = Pt(60)
p.font.color.rgb = SLATE

# Subtitle
sub_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(9), Inches(0.8))
tf = sub_box.text_frame
tf.text = "Subtitle goes here"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.color.rgb = SLATE_500

# Meta
meta_box = slide.shapes.add_textbox(Inches(2), Inches(4.8), Inches(9), Inches(0.6))
tf = meta_box.text_frame
tf.text = "Version v1.0.0 · 2025-11-01 · Authors: Your Name"
p = tf.paragraphs[0]
p.font.size = Pt(16)
p.font.color.rgb = SLATE_600

# ===== Section slide =====
slide = prs.slides.add_slide(slide_layout)
box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(10.5), Inches(1.5))
tf = box.text_frame
tf.text = "Section Title"
p = tf.paragraphs[0]
p.font.bold = True
p.font.size = Pt(60)
p.font.color.rgb = SLATE
p.alignment = PP_ALIGN.CENTER

# ===== Content slide =====
slide = prs.slides.add_slide(slide_layout)
hd = slide.shapes.add_textbox(Inches(1.5), Inches(1), Inches(10.5), Inches(1))
tf = hd.text_frame
tf.text = "Introduction"
p = tf.paragraphs[0]
p.font.bold = True
p.font.size = Pt(44)
p.font.color.rgb = SLATE

body = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10.5), Inches(3))
tf = body.text_frame
tf.text = "• Background and motivation\n• Design goals & constraints\n• Threat model overview"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.color.rgb = SLATE

# ===== Quote slide =====
slide = prs.slides.add_slide(slide_layout)
quote = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(2))
tf = quote.text_frame
tf.text = "“One idea per slide. Keep it crisp.”"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.color.rgb = SLATE
p.font.italic = True
p.alignment = PP_ALIGN.CENTER

prs.save("contrib/docs-templates/whitepaper.potx")
print("✅ Saved contrib/docs-templates/whitepaper.potx")
